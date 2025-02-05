# import : 다른 모듈이나 라이브러리의 코드를 사용할 수 있게 해주는 키워드
# sys모듈은 파이썬이 실행되는 환경과 소통하는 도구
# re모듈은 정규표현식을 사용할 수 있게 해주는 search(), match(), split(), findall().. 등등
import sys   
import json
import re
# PyQt5 라이브러리의 QtWidgets 모듈에서 특정 클래스들을 가져오는 것
# QApplication: PyQt5 애플리케이션의 기본 클래스로, 모든 GUI 프로그램에 필요합니다.
# QMainWindow: 메인 윈도우를 생성하는 클래스
# QLineEdit: 한 줄의 텍스트를 입력받는 위젯
# QPushButton: 클릭 가능한 버튼을 생성하는 클래스
# QTextEdit: 여러 줄의 텍스트를 입력하거나 표시할 수 있는 위젯
# QVBoxLayout: 위젯들을 수직으로 배열하는 레이아웃 클래스
# QWidget: 모든 UI 객체의 기본 클래스로, 다른 위젯들의 컨테이너 역할
# QComboBox: 드롭다운 목록을 제공하는 위젯
from PyQt5.QtWidgets import QApplication, QMainWindow, QLineEdit, QPushButton, QTextEdit, QVBoxLayout, QWidget, QComboBox
# python-pptx 라이브러리에서 Presentation 클래스를 가져오는 import 문
from pptx import Presentation
# python-pptx 라이브러리의 util 모듈에서 Inches 클래스를 가져오는 Inches 클래스는 길이를 인치 단위로 쉽게 지정할 수 있게 해주는 편리한 생성자
from pptx.util import Inches

# ✅ 성경 데이터 불러오기
#def 새로운 함수를 정의할 때 사용
def load_bible(file_path):   #load_bible 이라는 함수정의 file_path라는 매개변수를 받는다 
    with open(file_path, "r", encoding="utf-8-sig") as f: #with문은 파일을 자동으로 닫아줌,
        return json.load(f) #열린 파일 객체 'f'에서 json데이터를 파이썬 객체로 변환후 반환해줌줌

# ✅ 책 이름 약어 변환 딕셔너리 (전체)
# 딕셔너리 형태의 약어 설정 , {Key : value} 형태로 key를 입력하면 value가 호출됨 
book_abbreviations = {
    "창세기": "창", "창세": "창", "창": "창", "ㅊ": "창",
    "출애굽기": "출", "출애": "출", "출": "출", "ㅊㅇ": "출",
    "레위기": "레", "레위": "레", "ㄹㅇ": "레",
    "민수기": "민", "민수": "민", "ㅁㅅ": "민",
    "신명기": "신", "신명": "신", "ㅅㅁ": "신",
    "여호수아": "수", "여호": "수", "ㅇㅎ": "수",
    "사사기": "삿", "사사": "삿", "ㅅㅅ": "삿",
    "룻기": "룻", "ㄹ": "룻",
    "사무엘상": "삼상", "ㅅㅁㅇㅅ": "삼상",
    "사무엘하": "삼하", "ㅅㅁㅇㅎ": "삼하",
    "열왕기상": "왕상", "ㅇㅇㄱㅅ": "왕상",
    "열왕기하": "왕하", "ㅇㅇㄱㅎ": "왕하",
    "역대상": "대상", "ㅇㄷㅅ": "대상",
    "역대하": "대하", "ㅇㄷㅎ": "대하",
    "에스라": "스", "스라": "스", "ㅇㅅㄹ": "스",
    "느헤미야": "느", "느헤": "느", "ㄴㅎ": "느",
    "에스더": "에", "더": "에", "ㅇㅅㄷ": "에",
    "욥기": "욥", "ㅇ": "욥",
    "시편": "시", "ㅅㅍ": "시",
    "잠언": "잠", "ㅈㅇ": "잠",
    "전도서": "전", "전도": "전", "ㅈㄷ": "전",
    "아가": "아", "ㅇㄱ": "아",
    "이사야": "사", "사야": "사", "ㅇㅅㅇ": "사",
    "예레미야": "렘", "예레": "렘", "ㅇㄹㅁㅇ": "렘",
    "예레미야애가": "애", "애가": "애", "ㅇㄹㅁㅇㅇㄱ": "애",
    "에스겔": "겔", "ㅇㅅㄱ": "겔",
    "다니엘": "단", "니엘": "단", "ㄷㄴㅇ": "단",
    "호세아": "호", "호세": "호", "ㅎㅅㅇ": "호",
    "요엘": "욜", "엘": "욜", "ㅇㅇ": "욜",
    "아모스": "암", "모스": "암", "ㅇㅁㅅ": "암",
    "오바댜": "옵", "바댜": "옵", "ㅇㅂㄷ": "옵",
    "요나": "욘", "나": "욘", "ㅇㄴ": "욘",
    "미가": "미", "가": "미", "ㅁㄱ": "미",
    "나훔": "나", "훔": "나", "ㄴㅎ": "나",
    "하박국": "합", "박국": "합", "ㅎㅂㄱ": "합",
    "스바냐": "습", "바냐": "습", "ㅅㅂㄴ": "습",
    "학개": "학", "개": "학", "ㅎㄱ": "학",
    "스가랴": "슥", "가랴": "슥", "ㅅㄱㄹ": "슥",
    "말라기": "말", "라기": "말", "ㅁㄹㄱ": "말",
    "마태복음": "마", "마태": "마", "ㅁㅌ": "마",
    "마가복음": "막", "마가": "막", "ㅁㄱ": "막",
    "누가복음": "눅", "누가": "눅", "ㄴㄱ": "눅",
    "요한복음": "요", "요한": "요", "ㅇㅎ": "요",
    "사도행전": "행", "행전": "행", "ㅅㄷㅎㅈ": "행",
    "로마서": "롬", "로마": "롬", "ㄹㅁ": "롬",
    "고린도전서": "고전", "ㄱㄹㄷㅈ": "고전",
    "고린도후서": "고후", "ㄱㄹㄷㅎ": "고후",
    "갈라디아서": "갈", "갈라": "갈", "ㄱㄹㄷ": "갈",
    "에베소서": "엡", "에베": "엡", "ㅇㅂㅅ": "엡",
    "빌립보서": "빌", "빌립": "빌", "ㅂㄹㅂ": "빌",
    "골로새서": "골", "골로": "골", "ㄱㄹㅅ": "골",
    "데살로니가전서": "살전", "ㄷㅅㄹㄴㄱㅈ": "살전",
    "데살로니가후서": "살후", "ㄷㅅㄹㄴㄱㅎ": "살후",
    "디모데전서": "딤전", "ㄷㅁㄷㅈ": "딤전",
    "디모데후서": "딤후", "ㄷㅁㄷㅎ": "딤후",
    "디도서": "딛", "도서": "딛", "ㄷㄷ": "딛",
    "빌레몬서": "몬", "레몬": "몬", "ㅂㄹㅁ": "몬",
    "히브리서": "히", "히브": "히", "ㅎㅂㄹ": "히",
    "야고보서": "약", "고보": "약", "ㅇㄱㅂ": "약",
    "베드로전서": "벧전", "ㅂㄷㄹㅈ": "벧전",
    "베드로후서": "벧후", "ㅂㄷㄹㅎ": "벧후",
    "요한일서": "요일", "ㅇㅎㅇ": "요일",
    "요한이서": "요이", "ㅇㅎㅇ": "요이",
    "요한삼서": "요삼", "ㅇㅎㅅ": "요삼",
    "유다서": "유", "유다": "유", "ㅇㄷ": "유",
    "요한계시록": "계", "계시": "계", "ㅇㅎㄱㅅㄹ": "계"
}

# ✅ 검색어 변환 함수 
def parse_query(query):
    #re.sub 문자치열함수
    query = re.sub(r'\s+|장|절', '', query).replace(',', ':') # 주어진 문자열(auery)을 특정패턴(r'\s+|장|절')을 찾아아 다른 문자열('')로 대체
    pattern = r'^([가-힣ㄱ-ㅎ]+)?\s*(\d+):(\d+)(?:-(\d+))?$' # 검색어 분석을 위한 정규표현식의 패턴정의
    match = re.match(pattern, query) # pattern으로 정의된 정규표현식 패턴과 'query'문자열을 매칭시킴

    if not match:
        return None, None, None, None

    book, chapter, start_verse, end_verse = match.groups() # 정규 표현식의 각 괄호로 묶인 그룹에 해당하는 문자열들을 튜플로 반환
    
    if book: # book인 변수가 '참'일 때 실행, 문자열이 비어있지 않으면 '참'
        book = book_abbreviations.get(book, book) # get() 딕셔너리에서 키를 찾아 값을 반환, 첫번째인자는 찾고자 하는 '키', 두번째인자는 키가 없을 때 반환 시킬 '기본값'
    
    return book, chapter, start_verse, end_verse or start_verse


# ✅ PPT 파일 생성 함수
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_ppt(book, chapter, verse_range, text):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.62598425)
    
    # 본문 분할 (2절 단위)
    verses = text.split('\n')[2:]
    slide_groups = [verses[i:i+2] for i in range(0, len(verses), 2)]
    
    for group in slide_groups:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # 배경색 설정 (검정)
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = RGBColor(0, 0, 0)
        
        # 텍스트 상자 설정 (변경 부분)
        textbox = slide.shapes.add_textbox(
            left=Inches(0.1),  # 좌측 여백 증가
            top=Inches(0.15),   # 상단 여백 증가
            width=Inches(9.8),  # 폭 조정
            height=Inches(5.25984251)   # 높이 조정
        )
        tf = textbox.text_frame
        tf.word_wrap = True    # ✅ 줄바꿈 활성화
        tf.auto_size = None    # ✅ 자동 크기조정 비활성화
        
        # 제목 설정 (노란색)
        title = tf.add_paragraph()
        title.text = f"[{book} {chapter}:{verse_range}]"
        title.font.name = '맑은 고딕'
        title.font.size = Pt(27)
        title.font.color.rgb = RGBColor(255, 255, 0)
        title.font.bold = True  # 볼드 처리
        title.alignment = PP_ALIGN.LEFT
        
        # 본문 설정 (흰색)
        for verse in group:
            verse_text = verse.strip()
            if verse_text:
                p = tf.add_paragraph()
                # 절 번호와 내용을 분리하고 콜론을 제거합니다
                verse_parts = verse_text.split(':', 1)
                if len(verse_parts) == 2:
                    verse_num, verse_content = verse_parts
                    p.text = f"{verse_num} {verse_content.strip()}"
                else:
                    p.text = verse_text
                p.font.name = '맑은 고딕'
                p.font.size = Pt(27)
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.font.bold = True  # 볼드 처리
                p.alignment = PP_ALIGN.LEFT
                p.space_after = Pt(12)
    
    filename = f"{book}_{chapter}_{verse_range}.pptx"
    prs.save(filename)
    return filename


class BibleSearchApp(QMainWindow):
    def __init__(self):
        super().__init__()
        self.load_bible_data()
        self.initUI()

    def load_bible_data(self):
        self.bible_data = {
        "개역개정": load_bible("C:/Python/bible.json"),
        "새번역": load_bible("C:/Python/new_bible.json")
    }
            
    def initUI(self):
        self.setWindowTitle('성경 검색 프로그램')
        self.setGeometry(100, 100, 600, 400)

        central_widget = QWidget()
        self.setCentralWidget(central_widget)
        layout = QVBoxLayout(central_widget)

        self.version_select = QComboBox()
        self.version_select.addItems(["개역개정", "새번역"])
        self.search_input = QLineEdit()
        self.search_button = QPushButton('검색')
        self.result_text = QTextEdit()
        self.ppt_button = QPushButton('PPT 생성')

        layout.addWidget(self.version_select)
        layout.addWidget(self.search_input)
        layout.addWidget(self.search_button)
        layout.addWidget(self.result_text)
        layout.addWidget(self.ppt_button)

        self.search_button.clicked.connect(self.search_bible)
        self.ppt_button.clicked.connect(self.generate_ppt)
        self.search_input.returnPressed.connect(self.search_bible)

    def search_bible(self):
        search_query = self.search_input.text()
        selected_version = self.version_select.currentText()
        book, chapter, start_verse, end_verse = parse_query(search_query)

        if book and chapter and start_verse:
            start_key = f"{book}{chapter}:{start_verse}"
            end_key = f"{book}{chapter}:{end_verse}" if end_verse else start_key
            
            verses = {}
            for k, v in self.bible_data[selected_version].items():
                if k.startswith(f"{book}{chapter}:"):
                    verse_num = int(k.split(':')[1])
                    if int(start_verse) <= verse_num <= (int(end_verse) if end_verse else int(start_verse)):
                        verses[k] = v
            
            if verses:
                book_name = next((k for k, v in book_abbreviations.items() if v == book), book)
                result = f"{book_name} {chapter}장\n\n"
                
                for k, v in sorted(verses.items(), key=lambda x: int(x[0].split(':')[1])):
                    verse_num = k.split(':')[1]
                    result += f"{verse_num}: {v}\n"
                
                self.result_text.setText(result)
                self.ppt_file = create_ppt(book_name, chapter, f"{start_verse}-{end_verse}", result)
            else:
                self.result_text.setText("해당 구절을 찾을 수 없습니다.")
        else:
            self.result_text.setText("올바른 검색 형식이 아닙니다.")

    def generate_ppt(self):
        # PPT 생성 로직 구현 (이미 search_bible에서 생성됨)
        pass

if __name__ == '__main__':
    app = QApplication(sys.argv)
    ex = BibleSearchApp()
    ex.show()
    sys.exit(app.exec_())